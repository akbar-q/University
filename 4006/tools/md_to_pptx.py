"""Convert a Markdown slide deck (--- separators) into a .pptx.

Assumptions:
- Slides are separated by a line containing only: ---
- The first Markdown heading in a slide (e.g., # Title) becomes the slide title.
- Remaining non-empty lines become bullet points (basic conversion).

This is intentionally simple: it produces a usable PPTX skeleton students can edit in PowerPoint.

Usage:
  python tools/md_to_pptx.py templates/LO4_Fault_Finding_Presentation_Template.md out/LO4_Template.pptx

Dependencies:
  pip install python-pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation


_SLIDE_SEPARATOR_RE = re.compile(r"^\s*---\s*$", re.MULTILINE)
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)\s*$")
_LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")


def _strip_md(text: str) -> str:
    """Lightweight Markdown cleanup for slide text."""
    text = text.strip()
    if not text:
        return ""

    # Convert links to: text (url)
    def _link_sub(match: re.Match[str]) -> str:
        label = match.group(1).strip()
        url = match.group(2).strip()
        if not url:
            return label
        return f"{label} ({url})"

    text = _LINK_RE.sub(_link_sub, text)

    # Remove common emphasis/code markers.
    text = text.replace("**", "").replace("*", "").replace("`", "")

    # Collapse multiple spaces.
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _split_slides(markdown: str) -> list[str]:
    # Split on separator lines; keep non-empty chunks.
    parts = _SLIDE_SEPARATOR_RE.split(markdown)
    return [p.strip() for p in parts if p.strip()]


def _extract_title_and_body(slide_md: str) -> tuple[str, list[tuple[int, str]]]:
    """Return (title, bullets) where bullets are (level, text)."""
    lines = [ln.rstrip() for ln in slide_md.splitlines()]

    title = ""
    start_idx = 0

    for i, ln in enumerate(lines):
        m = _HEADING_RE.match(ln.strip())
        if m:
            title = _strip_md(m.group(2))
            start_idx = i + 1
            break

    bullets: list[tuple[int, str]] = []
    for ln in lines[start_idx:]:
        raw = ln.rstrip()
        if not raw.strip():
            continue
        if raw.strip().startswith("<!--"):
            continue

        # Determine indentation level (2 spaces or 1 tab per level).
        indent_spaces = len(raw) - len(raw.lstrip(" "))
        indent_tabs = len(raw) - len(raw.lstrip("\t"))
        level = 0
        if indent_tabs:
            level = indent_tabs
        elif indent_spaces:
            level = indent_spaces // 2

        content = raw.strip()

        # Remove common list prefixes.
        content = re.sub(r"^[-*+]\s+", "", content)
        content = re.sub(r"^\d+\.\s+", "", content)

        content = _strip_md(content)
        if content:
            bullets.append((min(level, 5), content))

    if not title:
        # Fallback: use first bullet as title if it looks like a title-ish line.
        if bullets:
            title = bullets[0][1][:80]
            bullets = bullets[1:]
        else:
            title = "(Untitled Slide)"

    return title, bullets


def build_pptx(input_md: Path, output_pptx: Path) -> None:
    markdown = input_md.read_text(encoding="utf-8")

    prs = Presentation()

    slides = _split_slides(markdown)
    for slide_md in slides:
        title, bullets = _extract_title_and_body(slide_md)

        # Layout 1 is typically "Title and Content".
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        body = None
        for shape in slide.placeholders:
            if shape.is_placeholder and shape.placeholder_format.type == 2:  # BODY
                body = shape
                break
        if body is None:
            # Fallback: first placeholder with text frame.
            for shape in slide.placeholders:
                if getattr(shape, "has_text_frame", False):
                    body = shape
                    break

        if body is not None and getattr(body, "has_text_frame", False):
            tf = body.text_frame
            tf.clear()

            for idx, (level, text) in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = text
                p.level = level

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert Markdown (--- separated slides) to PPTX")
    parser.add_argument("input", type=Path, help="Input Markdown file")
    parser.add_argument("output", type=Path, help="Output .pptx path")

    default_input = Path("templates/LO4_Fault_Finding_Presentation_Template.md")
    default_output = Path("out/LO4_Fault_Finding_Template.pptx")

    # Training-wheels behaviour: if run with no args, try sensible defaults.
    if len(sys.argv) == 1:
        if default_input.exists():
            args = parser.parse_args([str(default_input), str(default_output)])
            print(f"No arguments provided. Using default input: {default_input}")
            print(f"Writing default output: {default_output}")
        else:
            parser.print_help()
            print("\nTip: run with input and output paths, for example:")
            print("  python tools/md_to_pptx.py templates/LO4_Fault_Finding_Presentation_Template.md out/LO4_Fault_Finding_Template.pptx")
            return 2
    else:
        args = parser.parse_args()

    if not args.input.exists():
        raise SystemExit(f"Input file not found: {args.input}")
    if args.output.suffix.lower() != ".pptx":
        raise SystemExit("Output must end with .pptx")

    build_pptx(args.input, args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
